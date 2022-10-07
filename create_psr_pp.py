import collections 
import collections.abc
from pptx import Presentation
from datetime import date
from pptx.util import Inches
import os
import warnings

histogram_path=r'C:\Users\steineph\switchdrive\Private\Datenanalyse2\timedelta_histograms_zoomed_extra'
heatmap_path=r'C:\Users\steineph\switchdrive\Private\Datenanalyse4\PSR'

prs=Presentation()

#make it widescreen
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

lyt=prs.slide_layouts[0] # choosing a slide layout
slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1] # placeholder for subtitle
title.text="PSR UWO Sensoren" # title
subtitle.text="automatically generated on {}".format(date.today()) # subtitle

for hist_folder, heatmap_folder in zip(sorted(os.listdir(histogram_path)), sorted(os.listdir(heatmap_path))):
    if hist_folder==heatmap_folder:

        #add a slide defining which types of sensors we are looking at
        blank_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide_layout)
        title=slide.shapes.title # assigning a title
        title.text='Plots for {} sensors'.format(hist_folder.split('_')[0]) # title

        #iterate thorugh all plots and add them to the slide
        hist_image_path=os.path.join(histogram_path, hist_folder)
        heatmap_image_folder=os.path.join(heatmap_path, heatmap_folder)
        for img in sorted(os.listdir(heatmap_image_folder)):
            slide = prs.slides.add_slide(blank_slide_layout)
            title=slide.shapes.title # assigning a title
            title.text=img

            #positition of heatmap:
            left=Inches(0.5)
            top=Inches(1.2)
            img1=slide.shapes.add_picture(os.path.join(heatmap_image_folder, img),left,top)

            #positition of heatmap:
            left=Inches(9)
            top=Inches(0)
            img2=slide.shapes.add_picture(os.path.join(hist_image_path, img),left,top, height=Inches(10))
    
    
    else:
        warnings.warn('histfolder not the same as heatmapfolder')
prs.save("PSR_Histo_2.pptx")
        



